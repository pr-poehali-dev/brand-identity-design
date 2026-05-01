"""
Генерирует PowerPoint-презентацию бренда AURUM и загружает в S3.
Возвращает ссылку для скачивания файла.
"""
import os
import io
import json
import base64
import requests
import boto3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image

# Цветовая палитра AURUM
IVORY   = RGBColor(0xF5, 0xF0, 0xE8)
CREAM   = RGBColor(0xED, 0xE5, 0xD8)
SAND    = RGBColor(0xC8, 0xB8, 0x9A)
WARM_GRAY = RGBColor(0x8A, 0x7F, 0x74)
GRAPHITE  = RGBColor(0x1C, 0x19, 0x16)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# Изображения бренда
IMAGES = {
    "textile":      "https://cdn.poehali.dev/projects/3b373ac1-de12-4881-aa0b-a0a5a569a4b2/files/8c6054c2-ba45-4065-9b6c-cd5aea6ef9fa.jpg",
    "lookbook":     "https://cdn.poehali.dev/projects/3b373ac1-de12-4881-aa0b-a0a5a569a4b2/files/acf20eb2-8a72-4389-80d8-bce002f0d0bf.jpg",
    "campaign":     "https://cdn.poehali.dev/projects/3b373ac1-de12-4881-aa0b-a0a5a569a4b2/files/3230b108-afb8-4c98-a1d0-505b6a416b49.jpg",
    "flatlay":      "https://cdn.poehali.dev/projects/3b373ac1-de12-4881-aa0b-a0a5a569a4b2/files/25004912-cd7a-430c-9c80-acab1a251125.jpg",
    "model":        "https://cdn.poehali.dev/projects/3b373ac1-de12-4881-aa0b-a0a5a569a4b2/files/406c26ee-ac2b-402c-85dc-32126c1fa8a5.jpg",
    "web_mockup":   "https://cdn.poehali.dev/projects/3b373ac1-de12-4881-aa0b-a0a5a569a4b2/files/678f587b-16ad-4f16-afc9-5e63acf594ee.jpg",
    "fabrics":      "https://cdn.poehali.dev/projects/3b373ac1-de12-4881-aa0b-a0a5a569a4b2/files/b64200ec-edfa-4071-b64d-42e50e9ae5f5.jpg",
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def fetch_image(url: str, max_px: int = 1200) -> io.BytesIO:
    r = requests.get(url, timeout=15)
    r.raise_for_status()
    img = Image.open(io.BytesIO(r.content)).convert("RGB")
    if max(img.size) > max_px:
        img.thumbnail((max_px, max_px), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=75, optimize=True)
    buf.seek(0)
    return buf


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h,
             font_size=24, bold=False, italic=False,
             color=GRAPHITE, align=PP_ALIGN.LEFT,
             font_name="Montserrat"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, GRAPHITE)

    img = fetch_image(IMAGES["campaign"])
    pic = slide.shapes.add_picture(img, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    pic.zorder = 1

    # Затемнение
    overlay = add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=GRAPHITE)
    overlay.fill.fore_color.rgb = GRAPHITE
    from pptx.util import Pt as Pt2
    from pptx.dml.color import RGBColor as RGB2
    sp = overlay._element
    spPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = GRAPHITE
    overlay.fill.fore_color.theme_color  # touch
    # Задаём прозрачность через xml
    import lxml.etree as etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    solidFill = overlay._element.find('.//{%s}solidFill' % ns)
    if solidFill is not None:
        srgb = solidFill.find('{%s}srgbClr' % ns)
        if srgb is not None:
            alpha = etree.SubElement(srgb, '{%s}alpha' % ns)
            alpha.set('val', '75000')  # 75% opacity

    # Декоративный прямоугольник
    add_rect(slide, Inches(9.5), Inches(1), Inches(2.5), Inches(2.5),
             line_color=SAND, line_width=Pt(0.75))
    add_rect(slide, Inches(10), Inches(1.5), Inches(1.5), Inches(1.5),
             line_color=SAND, line_width=Pt(0.5))

    # Подзаголовок
    add_text(slide, "УСТОЙЧИВАЯ МОДА · АВАНГАРДНЫЙ ДИЗАЙН",
             Inches(0.8), Inches(3.2), Inches(8), Inches(0.5),
             font_size=10, color=SAND, font_name="Montserrat")

    # Логотип
    add_text(slide, "AURUM",
             Inches(0.8), Inches(3.7), Inches(10), Inches(2.2),
             font_size=120, bold=False, color=IVORY, font_name="Georgia")

    # Слоган
    add_text(slide, "Одежда, рождённая из земли — возвращающаяся к ней",
             Inches(0.8), Inches(6.0), Inches(7), Inches(0.8),
             font_size=14, italic=True, color=CREAM, font_name="Georgia")

    # Год
    add_text(slide, "Берлин · 2024",
             Inches(11), Inches(6.8), Inches(2), Inches(0.5),
             font_size=9, color=WARM_GRAY, align=PP_ALIGN.RIGHT, font_name="Montserrat")


def slide_positioning(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)

    # Номер слайда
    add_text(slide, "01", Inches(0.5), Inches(0.4), Inches(1), Inches(0.4),
             font_size=9, color=SAND, font_name="Montserrat")
    add_text(slide, "ПОЗИЦИОНИРОВАНИЕ БРЕНДА",
             Inches(1.2), Inches(0.4), Inches(8), Inches(0.4),
             font_size=9, color=WARM_GRAY, font_name="Montserrat")
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(12.3), Pt(1), fill_color=SAND)

    # Заголовок
    add_text(slide, "Мода как\nответственность",
             Inches(0.5), Inches(1.2), Inches(6), Inches(2),
             font_size=44, italic=True, color=GRAPHITE, font_name="Georgia")

    # Текстовые блоки
    texts = [
        ("БРЕНД", "AURUM — берлинский бренд устойчивой моды, основанный в 2022 году тремя дизайнерами с разных континентов."),
        ("МИССИЯ", "Доказать, что авангардная эстетика и экологическая ответственность не противоречат, а усиливают друг друга."),
        ("ЦЕЛЕВАЯ АУДИТОРИЯ", "Молодые профессионалы 22–35 лет, которые выбирают осознанно и хотят, чтобы их стиль отражал ценности."),
    ]
    y = Inches(3.3)
    for label, body in texts:
        add_text(slide, label, Inches(0.5), y, Inches(5.5), Inches(0.3),
                 font_size=8, color=SAND, font_name="Montserrat")
        add_text(slide, body, Inches(0.5), y + Inches(0.3), Inches(5.5), Inches(0.7),
                 font_size=12, color=GRAPHITE, font_name="Montserrat")
        y += Inches(1.1)

    # Ключевые сообщения справа
    add_rect(slide, Inches(7), Inches(1.2), Inches(5.8), Inches(5.8), fill_color=CREAM)
    add_text(slide, "КЛЮЧЕВЫЕ СООБЩЕНИЯ",
             Inches(7.3), Inches(1.5), Inches(5), Inches(0.4),
             font_size=8, color=WARM_GRAY, font_name="Montserrat")

    messages = [
        "«Красота не должна причинять вред»",
        "«Геометрия — наш язык»",
        "«Каждый шов — решение»",
        "«Земля вдохновляет, земле возвращаем»",
    ]
    my = Inches(2.1)
    for msg in messages:
        add_rect(slide, Inches(7.3), my, Pt(3), Inches(0.35), fill_color=SAND)
        add_text(slide, msg, Inches(7.7), my, Inches(4.8), Inches(0.45),
                 font_size=13, italic=True, color=GRAPHITE, font_name="Georgia")
        my += Inches(0.9)


def slide_product(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, GRAPHITE)

    add_text(slide, "02", Inches(0.5), Inches(0.4), Inches(1), Inches(0.4),
             font_size=9, color=SAND, font_name="Montserrat")
    add_text(slide, "КОЛЛЕКЦИЯ — TERRA FORMA SS24",
             Inches(1.2), Inches(0.4), Inches(8), Inches(0.4),
             font_size=9, color=WARM_GRAY, font_name="Montserrat")
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(12.3), Pt(1), fill_color=SAND)

    # Большое фото — лукбук
    img_lb = fetch_image(IMAGES["lookbook"])
    slide.shapes.add_picture(img_lb, Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.6))

    # Маленькое фото — текстиль
    img_tx = fetch_image(IMAGES["textile"])
    slide.shapes.add_picture(img_tx, Inches(8.2), Inches(1.1), Inches(4.6), Inches(2.7))

    # Описание
    add_rect(slide, Inches(8.2), Inches(3.9), Inches(4.6), Inches(2.8), fill_color=CREAM)
    add_text(slide, "Terra Forma",
             Inches(8.4), Inches(4.1), Inches(4.2), Inches(0.8),
             font_size=28, italic=True, color=GRAPHITE, font_name="Georgia")
    add_text(slide, "Земля как основа — геометрия как язык.\nКоллекция исследует архетипы формы\nчерез переработанные материалы.",
             Inches(8.4), Inches(4.9), Inches(4.2), Inches(1.5),
             font_size=11, color=GRAPHITE, font_name="Montserrat")

    add_text(slide, "Весна — Лето 2024",
             Inches(0.5), Inches(6.9), Inches(4), Inches(0.4),
             font_size=9, color=SAND, font_name="Montserrat")


def slide_promo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)

    add_text(slide, "03", Inches(0.5), Inches(0.4), Inches(1), Inches(0.4),
             font_size=9, color=SAND, font_name="Montserrat")
    add_text(slide, "ПРОМО-МАТЕРИАЛЫ",
             Inches(1.2), Inches(0.4), Inches(8), Inches(0.4),
             font_size=9, color=WARM_GRAY, font_name="Montserrat")
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(12.3), Pt(1), fill_color=SAND)

    # Постер 1
    img_cp = fetch_image(IMAGES["campaign"])
    slide.shapes.add_picture(img_cp, Inches(0.5), Inches(1.1), Inches(4.0), Inches(5.3))

    # Постер 2
    img_lb = fetch_image(IMAGES["lookbook"])
    slide.shapes.add_picture(img_lb, Inches(4.7), Inches(1.1), Inches(4.0), Inches(5.3))

    # Баннер
    img_tx = fetch_image(IMAGES["textile"])
    slide.shapes.add_picture(img_tx, Inches(8.9), Inches(1.1), Inches(3.9), Inches(5.3))

    # Подписи
    captions = [
        (Inches(0.5), '"Форма как смысл"'),
        (Inches(4.7), '"Земля и тело"'),
        (Inches(8.9), '"Носи землю легко"'),
    ]
    for x, cap in captions:
        add_text(slide, cap, x, Inches(6.5), Inches(3.8), Inches(0.5),
                 font_size=12, italic=True, color=GRAPHITE, font_name="Georgia")

    add_text(slide, "AURUM — Terra Forma Collection — SS24",
             Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
             font_size=8, color=WARM_GRAY, align=PP_ALIGN.CENTER, font_name="Montserrat")


def slide_identity(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, GRAPHITE)

    add_text(slide, "04", Inches(0.5), Inches(0.4), Inches(1), Inches(0.4),
             font_size=9, color=SAND, font_name="Montserrat")
    add_text(slide, "АЙДЕНТИКА",
             Inches(1.2), Inches(0.4), Inches(8), Inches(0.4),
             font_size=9, color=WARM_GRAY, font_name="Montserrat")
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(12.3), Pt(1), fill_color=SAND)

    # Логотип крупно
    add_rect(slide, Inches(0.5), Inches(1.1), Inches(5.8), Inches(4), fill_color=IVORY)
    add_text(slide, "AURUM",
             Inches(0.7), Inches(1.5), Inches(5.4), Inches(2.5),
             font_size=72, color=GRAPHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
    add_text(slide, "SUSTAINABLE AVANT-GARDE FASHION",
             Inches(0.7), Inches(3.8), Inches(5.4), Inches(0.5),
             font_size=8, color=WARM_GRAY, align=PP_ALIGN.CENTER, font_name="Montserrat")

    # Инвертированная версия
    add_rect(slide, Inches(0.5), Inches(5.2), Inches(2.7), Inches(1.5), fill_color=GRAPHITE)
    add_rect(slide, Inches(0.5), Inches(5.2), Inches(2.7), Inches(1.5), line_color=SAND, line_width=Pt(0.5))
    add_text(slide, "AURUM", Inches(0.5), Inches(5.35), Inches(2.7), Inches(1.0),
             font_size=36, color=IVORY, align=PP_ALIGN.CENTER, font_name="Georgia")

    # Версия на песке
    add_rect(slide, Inches(3.4), Inches(5.2), Inches(2.9), Inches(1.5), fill_color=SAND)
    add_text(slide, "AURUM", Inches(3.4), Inches(5.35), Inches(2.9), Inches(1.0),
             font_size=36, color=GRAPHITE, align=PP_ALIGN.CENTER, font_name="Georgia")

    # Цветовая палитра справа
    colors_data = [
        (IVORY,     "IVORY 001",     "#F5F0E8", GRAPHITE),
        (CREAM,     "CREAM 002",     "#EDE5D8", GRAPHITE),
        (SAND,      "SAND 003",      "#C8B89A", GRAPHITE),
        (WARM_GRAY, "WARM GREY 004", "#8A7F74", IVORY),
        (GRAPHITE,  "GRAPHITE 005",  "#1C1916", IVORY),
    ]
    cx = Inches(6.6)
    add_text(slide, "ЦВЕТОВАЯ ПАЛИТРА",
             cx, Inches(1.1), Inches(6.5), Inches(0.4),
             font_size=8, color=WARM_GRAY, font_name="Montserrat")
    cy = Inches(1.6)
    for c, name, hexv, text_c in colors_data:
        add_rect(slide, cx, cy, Inches(6.5), Inches(0.85), fill_color=c)
        add_text(slide, name, cx + Inches(0.2), cy + Inches(0.2), Inches(4), Inches(0.5),
                 font_size=11, color=text_c, font_name="Montserrat")
        add_text(slide, hexv, cx + Inches(5.5), cy + Inches(0.2), Inches(1.2), Inches(0.5),
                 font_size=9, color=text_c, align=PP_ALIGN.RIGHT, font_name="Montserrat")
        cy += Inches(0.88)

    # Типографика
    add_text(slide, "ТИПОГРАФИКА",
             Inches(6.6), Inches(6.4), Inches(6), Inches(0.3),
             font_size=8, color=WARM_GRAY, font_name="Montserrat")
    add_text(slide, "Cormorant Garamond",
             Inches(6.6), Inches(6.7), Inches(4), Inches(0.6),
             font_size=20, italic=True, color=IVORY, font_name="Georgia")
    add_text(slide, "Montserrat Light / Regular",
             Inches(10), Inches(6.7), Inches(3.2), Inches(0.6),
             font_size=13, color=SAND, font_name="Montserrat")


def slide_application(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    add_text(slide, "05", Inches(0.5), Inches(0.4), Inches(1), Inches(0.4),
             font_size=9, color=SAND, font_name="Montserrat")
    add_text(slide, "ПРИМЕНЕНИЕ НА НОСИТЕЛЯХ",
             Inches(1.2), Inches(0.4), Inches(8), Inches(0.4),
             font_size=9, color=WARM_GRAY, font_name="Montserrat")
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(12.3), Pt(1), fill_color=SAND)

    add_text(slide, "Айдентика в контексте",
             Inches(0.5), Inches(1.1), Inches(8), Inches(0.8),
             font_size=36, italic=True, color=GRAPHITE, font_name="Georgia")

    # Носители
    carriers = [
        ("Упаковка", "Крафт-бумага с тиснением логотипа и цветовым акцентом SAND.\nМинималистичная, 100% переработанная."),
        ("Бирки", "Лаконичная бирка из переработанного картона с логотипом AURUM,\nсоставом ткани и QR-кодом экотрекинга."),
        ("Визитка", "Монохромный дизайн на слоновой кости, шрифт Cormorant,\nрельефное тиснение логотипа без цвета."),
        ("Шоппер", "Натуральный хлопок цвета IVORY, принт AURUM в графите.\nМногоразовый, биоразлагаемые чернила."),
        ("Соцсети", "Квадратные форматы 1:1 и 4:5, строгая сетка публикаций,\nнейтральный фон, минимум текста."),
        ("Рассылка", "Текст на молочном фоне, Cormorant для заголовков,\nMontserrat для body, без агрессивных CTA."),
    ]
    cols = 3
    cw = Inches(4.0)
    ch = Inches(1.8)
    for i, (title, desc) in enumerate(carriers):
        col = i % cols
        row = i // cols
        cx = Inches(0.5) + col * Inches(4.3)
        cy = Inches(2.2) + row * Inches(2.1)
        add_rect(slide, cx, cy, cw, ch, fill_color=IVORY)
        add_rect(slide, cx, cy, Pt(4), ch, fill_color=SAND)
        add_text(slide, title, cx + Inches(0.15), cy + Inches(0.12), cw - Inches(0.2), Inches(0.45),
                 font_size=14, bold=False, color=GRAPHITE, font_name="Georgia")
        add_text(slide, desc, cx + Inches(0.15), cy + Inches(0.55), cw - Inches(0.2), Inches(1.1),
                 font_size=9, color=WARM_GRAY, font_name="Montserrat")


def slide_collection_visuals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, GRAPHITE)

    add_text(slide, "02b", Inches(0.5), Inches(0.4), Inches(1), Inches(0.4),
             font_size=9, color=SAND, font_name="Montserrat")
    add_text(slide, "ОДЕЖДА КОЛЛЕКЦИИ — TERRA FORMA",
             Inches(1.2), Inches(0.4), Inches(9), Inches(0.4),
             font_size=9, color=WARM_GRAY, font_name="Montserrat")
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(12.3), Pt(1), fill_color=SAND)

    # Модель — большое фото слева
    img_model = fetch_image(IMAGES["model"])
    slide.shapes.add_picture(img_model, Inches(0.5), Inches(1.1), Inches(5.5), Inches(5.6))

    # Flat-lay справа вверху
    img_flat = fetch_image(IMAGES["flatlay"])
    slide.shapes.add_picture(img_flat, Inches(6.2), Inches(1.1), Inches(6.6), Inches(2.7))

    # Детали тканей справа внизу
    img_fab = fetch_image(IMAGES["fabrics"])
    slide.shapes.add_picture(img_fab, Inches(6.2), Inches(3.9), Inches(6.6), Inches(2.8))

    # Подписи
    add_text(slide, "Lookbook · Модель в пальто Terra Forma",
             Inches(0.5), Inches(6.8), Inches(5.5), Inches(0.4),
             font_size=9, italic=True, color=SAND, font_name="Montserrat")
    add_text(slide, "Flat-lay · Ключевые силуэты SS24",
             Inches(6.2), Inches(6.8), Inches(6), Inches(0.4),
             font_size=9, italic=True, color=SAND, font_name="Montserrat")


def slide_logo_rationale(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)

    add_text(slide, "04b", Inches(0.5), Inches(0.4), Inches(1), Inches(0.4),
             font_size=9, color=SAND, font_name="Montserrat")
    add_text(slide, "ЛОГИКА ЛОГОТИПА",
             Inches(1.2), Inches(0.4), Inches(8), Inches(0.4),
             font_size=9, color=WARM_GRAY, font_name="Montserrat")
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(12.3), Pt(1), fill_color=SAND)

    add_text(slide, "Почему AURUM выглядит именно так",
             Inches(0.5), Inches(1.1), Inches(8), Inches(0.9),
             font_size=36, italic=True, color=GRAPHITE, font_name="Georgia")

    # Демонстрация логотипа
    add_rect(slide, Inches(0.5), Inches(2.1), Inches(5.8), Inches(2.0), fill_color=GRAPHITE)
    add_text(slide, "AURUM",
             Inches(0.5), Inches(2.2), Inches(5.8), Inches(1.6),
             font_size=72, color=IVORY, align=PP_ALIGN.CENTER, font_name="Georgia")

    # Пояснения — три принципа
    principles = [
        ("Только текст, без символа",
         "Имя бренда — само является символом. Золото (aurum) в химии не нуждается в иконке — только в чистоте написания."),
        ("Cormorant Garamond",
         "Засечный шрифт с историей. Отсылка к традиции ручного шитья и высокой моды — при этом достаточно геометричный для авангарда."),
        ("Только заглавные буквы",
         "Равнозначность каждой буквы. Нет иерархии — как в устойчивой моде нет «главного» и «второстепенного» в цепочке производства."),
        ("Почему Берлин",
         "Берлин — город, где авангард и осознанность сосуществуют в культуре. Это не маркетинг, а точное попадание в среду бренда."),
    ]

    py = Inches(2.1)
    for i, (title, desc) in enumerate(principles):
        col = i % 2
        row = i // 2
        cx = Inches(6.6) if col == 0 else Inches(10.0)
        cy = Inches(2.1) + row * Inches(2.4)
        add_rect(slide, cx, cy, Inches(3.1), Inches(2.1), fill_color=CREAM)
        add_rect(slide, cx, cy, Inches(3.1), Pt(3), fill_color=SAND)
        add_text(slide, title, cx + Inches(0.15), cy + Inches(0.18), Inches(2.8), Inches(0.6),
                 font_size=11, bold=False, color=GRAPHITE, font_name="Georgia")
        add_text(slide, desc, cx + Inches(0.15), cy + Inches(0.7), Inches(2.8), Inches(1.2),
                 font_size=9, color=WARM_GRAY, font_name="Montserrat")

    # Промо-каналы для материалов
    add_rect(slide, Inches(0.5), Inches(4.3), Inches(5.8), Inches(1.0), fill_color=SAND)
    add_text(slide, "ГДЕ ИСПОЛЬЗУЮТСЯ ПРОМО-МАТЕРИАЛЫ",
             Inches(0.7), Inches(4.4), Inches(5.4), Inches(0.3),
             font_size=8, color=GRAPHITE, font_name="Montserrat")
    channels = [("Instagram / Reels", "1:1 и 4:5"), ("OOH / Афиши", "A1 формат"), ("Email-рассылка", "600px wide")]
    for i, (ch, fmt) in enumerate(channels):
        add_text(slide, f"· {ch}  —  {fmt}",
                 Inches(0.7), Inches(4.8) + i * Inches(0.22), Inches(5.4), Inches(0.25),
                 font_size=9, color=GRAPHITE, font_name="Montserrat")


def slide_web_mockup(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, CREAM)

    add_text(slide, "05b", Inches(0.5), Inches(0.4), Inches(1), Inches(0.4),
             font_size=9, color=SAND, font_name="Montserrat")
    add_text(slide, "DIGITAL-НОСИТЕЛЬ — ВЕБ-САЙТ",
             Inches(1.2), Inches(0.4), Inches(9), Inches(0.4),
             font_size=9, color=WARM_GRAY, font_name="Montserrat")
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(12.3), Pt(1), fill_color=SAND)

    add_text(slide, "Айдентика в digital-среде",
             Inches(0.5), Inches(1.1), Inches(6), Inches(0.8),
             font_size=36, italic=True, color=GRAPHITE, font_name="Georgia")

    # Мокап сайта — большое изображение
    img_web = fetch_image(IMAGES["web_mockup"])
    slide.shapes.add_picture(img_web, Inches(0.5), Inches(2.0), Inches(8.5), Inches(4.8))

    # Описание справа
    add_rect(slide, Inches(9.2), Inches(2.0), Inches(3.9), Inches(4.8), fill_color=GRAPHITE)

    web_points = [
        ("Цвет фона", "Graphite #1C1916"),
        ("Заголовки", "Cormorant Garamond Italic"),
        ("Текст", "Montserrat Light"),
        ("Акценты", "Sand #C8B89A"),
        ("Сетка", "12-колоночная, отступы 80px"),
        ("Мобильная версия", "Адаптирована 320px+"),
    ]
    add_text(slide, "ПАРАМЕТРЫ САЙТА",
             Inches(9.4), Inches(2.15), Inches(3.5), Inches(0.35),
             font_size=8, color=SAND, font_name="Montserrat")
    wy = Inches(2.6)
    for label, val in web_points:
        add_text(slide, label, Inches(9.4), wy, Inches(3.5), Inches(0.28),
                 font_size=8, color=WARM_GRAY, font_name="Montserrat")
        add_text(slide, val, Inches(9.4), wy + Inches(0.27), Inches(3.5), Inches(0.32),
                 font_size=11, color=IVORY, font_name="Georgia")
        wy += Inches(0.62)

    add_text(slide, "aurum.berlin",
             Inches(0.5), Inches(6.9), Inches(4), Inches(0.4),
             font_size=10, color=WARM_GRAY, font_name="Montserrat")


def slide_ai_iterations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)

    add_text(slide, "06b", Inches(0.5), Inches(0.4), Inches(1), Inches(0.4),
             font_size=9, color=SAND, font_name="Montserrat")
    add_text(slide, "КАК РАБОТАЕТ ГЕНЕРАЦИЯ ВИЗУАЛОВ",
             Inches(1.2), Inches(0.4), Inches(9), Inches(0.4),
             font_size=9, color=WARM_GRAY, font_name="Montserrat")
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(12.3), Pt(1), fill_color=SAND)

    add_text(slide, "От промпта к результату",
             Inches(0.5), Inches(1.1), Inches(8), Inches(0.8),
             font_size=36, italic=True, color=GRAPHITE, font_name="Georgia")

    # Итерации — 3 шага
    iterations = [
        ("Итерация 1 — Концепция",
         '"sustainable fashion editorial, minimalist, ivory tones"',
         "Слишком нейтрально. Нет характера бренда, нет земляных оттенков."),
        ("Итерация 2 — Уточнение",
         '"avant-garde structured garments, clay sand beige, geometric folds, Berlin aesthetic"',
         "Появилась геометрия и палитра. Но не хватало ощущения «земли»."),
        ("Итерация 3 — Финал",
         '"AURUM Terra Forma, recycled linen coat, clay-dyed fabric, earthy minimalism, editorial photography"',
         "Точное попадание: бренд, материал, настроение — всё считывается с первого взгляда."),
    ]
    iy = Inches(2.1)
    for i, (title, prompt, result) in enumerate(iterations):
        bg = CREAM if i % 2 == 0 else IVORY
        add_rect(slide, Inches(0.5), iy, Inches(12.3), Inches(1.55), fill_color=bg)
        # Номер
        add_rect(slide, Inches(0.5), iy, Inches(0.4), Inches(1.55), fill_color=SAND if i == 2 else WARM_GRAY)
        add_text(slide, str(i + 1), Inches(0.5), iy + Inches(0.5), Inches(0.4), Inches(0.6),
                 font_size=14, bold=True, color=IVORY, align=PP_ALIGN.CENTER, font_name="Montserrat")
        # Заголовок
        add_text(slide, title, Inches(1.1), iy + Inches(0.1), Inches(11.5), Inches(0.4),
                 font_size=11, color=GRAPHITE, font_name="Georgia")
        # Промпт
        add_text(slide, prompt, Inches(1.1), iy + Inches(0.45), Inches(11.5), Inches(0.5),
                 font_size=9, italic=True, color=WARM_GRAY, font_name="Montserrat")
        # Результат
        marker = "✓" if i == 2 else "→"
        add_text(slide, f"{marker} {result}", Inches(1.1), iy + Inches(0.92), Inches(11.5), Inches(0.5),
                 font_size=10, color=GRAPHITE if i == 2 else WARM_GRAY, font_name="Montserrat")
        iy += Inches(1.65)

    # Вывод
    add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35), fill_color=GRAPHITE)
    add_text(slide, "Каждый визуал прошёл минимум 3 итерации уточнения промпта перед финальным выбором",
             Inches(0.7), Inches(7.04), Inches(12), Inches(0.28),
             font_size=9, color=SAND, font_name="Montserrat")


def slide_process(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, IVORY)

    add_text(slide, "07", Inches(0.5), Inches(0.4), Inches(1), Inches(0.4),
             font_size=9, color=SAND, font_name="Montserrat")
    add_text(slide, "ПРОЦЕСС СОЗДАНИЯ",
             Inches(1.2), Inches(0.4), Inches(8), Inches(0.4),
             font_size=9, color=WARM_GRAY, font_name="Montserrat")
    add_rect(slide, Inches(0.5), Inches(0.9), Inches(12.3), Pt(1), fill_color=SAND)

    add_text(slide, "Как создавался AURUM",
             Inches(0.5), Inches(1.1), Inches(7), Inches(0.9),
             font_size=36, italic=True, color=GRAPHITE, font_name="Georgia")

    # Инструменты — левая колонка
    tools = [
        ("Poehali.dev",        "Платформа создания"),
        ("Claude AI",          "Архитектура бренда и текст"),
        ("FLUX Image AI",      "Генерация визуалов"),
        ("React + TypeScript", "Frontend разработка"),
        ("Cormorant Garamond", "Дисплейная типографика"),
        ("python-pptx",        "Генерация этого файла"),
    ]
    add_text(slide, "ИНСТРУМЕНТЫ",
             Inches(0.5), Inches(2.2), Inches(5), Inches(0.35),
             font_size=8, color=SAND, font_name="Montserrat")
    ty = Inches(2.65)
    for i, (tool, role) in enumerate(tools):
        add_rect(slide, Inches(0.5), ty, Pt(3), Inches(0.42), fill_color=SAND)
        add_text(slide, f"{tool}", Inches(0.8), ty, Inches(2.5), Inches(0.42),
                 font_size=12, color=GRAPHITE, font_name="Georgia")
        add_text(slide, role, Inches(3.5), ty, Inches(3), Inches(0.42),
                 font_size=9, color=WARM_GRAY, font_name="Montserrat")
        ty += Inches(0.52)

    # Промпты — правая колонка
    add_rect(slide, Inches(7), Inches(1.1), Inches(5.8), Inches(6.1), fill_color=CREAM)
    add_text(slide, "КЛЮЧЕВЫЕ ПРОМПТЫ",
             Inches(7.2), Inches(1.3), Inches(5.4), Inches(0.35),
             font_size=8, color=WARM_GRAY, font_name="Montserrat")

    prompts = [
        ("Концепция", '"Бренд устойчивой моды с авангардным дизайном для молодого поколения. Чистые линии, нейтральная палитра, геометрические формы."'),
        ("Визуал", '"Minimalist sustainable fashion editorial, geometric folded fabric in ivory and warm sand tones, luxury fashion photography, soft natural light."'),
        ("Сайт", '"Создай сайт-презентацию бренда: hero, манифест, коллекция, материалы, палитра, постеры, процесс."'),
    ]
    py = Inches(1.8)
    for label, prompt in prompts:
        add_text(slide, label.upper(), Inches(7.2), py, Inches(5.4), Inches(0.3),
                 font_size=7, color=SAND, font_name="Montserrat")
        add_text(slide, prompt, Inches(7.2), py + Inches(0.3), Inches(5.4), Inches(1.0),
                 font_size=10, italic=True, color=GRAPHITE, font_name="Georgia")
        py += Inches(1.5)

    # Время создания
    add_rect(slide, Inches(0.5), Inches(6.2), Inches(5.5), Inches(1.0), fill_color=GRAPHITE)
    add_text(slide, "~15 минут", Inches(0.7), Inches(6.25), Inches(2.5), Inches(0.6),
             font_size=28, color=IVORY, font_name="Georgia")
    add_text(slide, "от идеи до полноценной презентации бренда",
             Inches(3.3), Inches(6.4), Inches(2.6), Inches(0.5),
             font_size=9, color=SAND, font_name="Montserrat")


def slide_end(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, GRAPHITE)

    img = fetch_image(IMAGES["textile"])
    slide.shapes.add_picture(img, Inches(6.5), Inches(0), Inches(6.83), SLIDE_H)

    add_rect(slide, Inches(6.5), Inches(0), Inches(6.83), SLIDE_H, fill_color=GRAPHITE)
    import lxml.etree as etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Геометрия
    add_rect(slide, Inches(9), Inches(1), Inches(2.5), Inches(2.5), line_color=SAND, line_width=Pt(0.75))

    add_text(slide, "AURUM",
             Inches(0.5), Inches(2.0), Inches(6), Inches(2.5),
             font_size=80, color=IVORY, font_name="Georgia")
    add_text(slide, "SUSTAINABLE AVANT-GARDE FASHION",
             Inches(0.5), Inches(4.5), Inches(6), Inches(0.5),
             font_size=9, color=SAND, font_name="Montserrat")

    add_rect(slide, Inches(0.5), Inches(5.2), Inches(3), Pt(1), fill_color=SAND)

    add_text(slide, "© 2024 AURUM — Концептуальный бренд-проект\nСоздано на poehali.dev",
             Inches(0.5), Inches(5.5), Inches(6), Inches(0.8),
             font_size=9, color=WARM_GRAY, font_name="Montserrat")


def build_pptx() -> bytes:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_positioning(prs)
    slide_product(prs)
    slide_collection_visuals(prs)
    slide_promo(prs)
    slide_identity(prs)
    slide_logo_rationale(prs)
    slide_application(prs)
    slide_web_mockup(prs)
    slide_process(prs)
    slide_ai_iterations(prs)
    slide_end(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


FILENAME = 'AURUM_Brand_Presentation_v2.pptx'
S3_KEY = f"pptx/{FILENAME}"


def get_s3_client():
    return boto3.client(
        's3',
        endpoint_url='https://bucket.poehali.dev',
        aws_access_key_id=os.environ['AWS_ACCESS_KEY_ID'],
        aws_secret_access_key=os.environ['AWS_SECRET_ACCESS_KEY'],
    )


def file_exists_in_s3() -> bool:
    s3 = get_s3_client()
    try:
        s3.head_object(Bucket='files', Key=S3_KEY)
        return True
    except Exception:
        return False


def upload_to_s3(data: bytes) -> str:
    s3 = get_s3_client()
    s3.put_object(
        Bucket='files',
        Key=S3_KEY,
        Body=data,
        ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        ContentDisposition=f'attachment; filename="{FILENAME}"',
    )
    return f"https://cdn.poehali.dev/projects/{os.environ['AWS_ACCESS_KEY_ID']}/bucket/{S3_KEY}"


def get_cdn_url() -> str:
    return f"https://cdn.poehali.dev/projects/{os.environ['AWS_ACCESS_KEY_ID']}/bucket/{S3_KEY}"


def list_all_files() -> list:
    s3 = get_s3_client()
    result = []
    try:
        paginator = s3.get_paginator('list_objects_v2')
        for page in paginator.paginate(Bucket='files'):
            result.extend(page.get('Contents', []))
    except Exception:
        pass
    return result


def cleanup_old_files():
    s3 = get_s3_client()
    try:
        all_files = list_all_files()
        deleted = []
        for obj in all_files:
            key = obj['Key']
            if key != S3_KEY:
                s3.delete_object(Bucket='files', Key=key)
                deleted.append(key)
        return deleted
    except Exception:
        return []


def handler(event: dict, context) -> dict:
    """Генерирует PPTX-презентацию AURUM и возвращает файл напрямую в base64."""
    cors = {
        'Access-Control-Allow-Origin': '*',
        'Access-Control-Allow-Methods': 'GET, POST, OPTIONS',
        'Access-Control-Allow-Headers': 'Content-Type',
    }

    if event.get('httpMethod') == 'OPTIONS':
        return {'statusCode': 200, 'headers': cors, 'body': ''}

    pptx_bytes = build_pptx()
    b64 = base64.b64encode(pptx_bytes).decode('utf-8')

    return {
        'statusCode': 200,
        'headers': {
            **cors,
            'Content-Type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'Content-Disposition': f'attachment; filename="{FILENAME}"',
        },
        'body': b64,
        'isBase64Encoded': True,
    }